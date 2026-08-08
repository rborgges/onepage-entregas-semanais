#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Preenche o onepage semanal de entregas (.pptx) a partir de um arquivo de
texto simples (ex.: editado no Notepad++).

USO:
    python preencher_onepage.py dados_semana.txt
    python preencher_onepage.py dados_semana.txt --template Onepage_Entregas_Semanais.pptx --output saida.pptx

FORMATO DO ARQUIVO .txt (ver dados_semana_exemplo.txt):

    EQUIPE: Time de Growth
    PERIODO: 04/08 a 08/08/2026
    RESPONSAVEL: Fulano de Tal

    ENTREGAS:
    Landing page nova | Ana | Concluído | No ar desde quarta
    Integração com API de pagamento | Bruno | Em andamento | 80% concluído
    Ajuste no checkout | Carla | Atrasado | Aguardando fornecedor

    DESTAQUES:
    Lançamento da campanha X
    Meta de conversão batida

    RISCOS:
    Atraso do fornecedor Z
    Falta de recurso pra sprint 2

Regras:
- Cada entrega é uma linha dentro de "ENTREGAS:", com 4 campos separados por "|":
  Nome da entrega | Responsável | Status | Observação (observação é opcional)
- Status reconhecidos (não sensível a maiúsculas/acentos): Concluído, Em andamento,
  Atrasado, Bloqueado. Qualquer outro texto entra sem cor especial.
- DESTAQUES e RISCOS aceitam qualquer número de linhas (uma por bullet).
- Uma seção em branco (sem linhas abaixo dela) é simplesmente ignorada.
"""

import argparse
import copy
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

# ---------------------------------------------------------------------------
# Cores (mesma paleta do template)
# ---------------------------------------------------------------------------
GREEN = RGBColor(0x1A, 0x87, 0x54)
GREEN_BG = RGBColor(0xE6, 0xF4, 0xEA)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG = RGBColor(0xFD, 0xF1, 0xDC)
RED = RGBColor(0xB4, 0x23, 0x18)
RED_BG = RGBColor(0xFB, 0xE8, 0xE6)
GRAY = RGBColor(0x37, 0x41, 0x51)
GRAY_BG = RGBColor(0xF0, 0xF1, 0xF3)

# Índices fixos dos shapes no template gerado (ver build.js)
IDX_TITULO = 0
IDX_EQUIPE = 1
IDX_PERIODO = 2
IDX_RESPONSAVEL = 3
IDX_TABELA = 11
IDX_DESTAQUES = 14
IDX_RISCOS = 17


def normalize(s: str) -> str:
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode("ascii")
    return s.strip().lower()


def status_colors(status: str):
    n = normalize(status)
    if n in ("concluido", "concluida", "feito", "done"):
        return GREEN, GREEN_BG
    if n in ("em andamento", "andamento", "in progress"):
        return AMBER, AMBER_BG
    if n in ("atrasado", "atrasada", "bloqueado", "bloqueada", "delayed", "blocked"):
        return RED, RED_BG
    return GRAY, GRAY_BG


# ---------------------------------------------------------------------------
# Parser do .txt
# ---------------------------------------------------------------------------
def read_text_file(path: Path) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def parse_data(text: str) -> dict:
    data = {
        "equipe": "",
        "periodo": "",
        "responsavel": "",
        "entregas": [],
        "destaques": [],
        "riscos": [],
    }
    section = None
    field_map = {
        "equipe": "equipe",
        "periodo": "periodo",
        "período": "periodo",
        "responsavel": "responsavel",
        "responsável": "responsavel",
    }
    section_map = {"entregas": "entregas", "destaques": "destaques", "riscos": "riscos"}

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        if ":" in line:
            key, _, rest = line.partition(":")
            key_norm = normalize(key)
            rest = rest.strip()

            if key_norm in section_map:
                section = section_map[key_norm]
                continue

            if key_norm in field_map and not rest == "" or (key_norm in field_map and section is None):
                # campo simples tipo "EQUIPE: valor"
                if key_norm in field_map:
                    data[field_map[key_norm]] = rest
                    section = None
                    continue

        if section == "entregas":
            parts = [p.strip() for p in line.split("|")]
            while len(parts) < 4:
                parts.append("")
            nome, responsavel, status, obs = parts[:4]
            data["entregas"].append(
                {"nome": nome, "responsavel": responsavel, "status": status, "obs": obs}
            )
        elif section == "destaques":
            data["destaques"].append(line)
        elif section == "riscos":
            data["riscos"].append(line)

    return data


# ---------------------------------------------------------------------------
# Helpers de preenchimento do pptx
# ---------------------------------------------------------------------------
def set_single_run_text(shape, new_text: str):
    """Troca o texto mantendo a formatação do único run existente."""
    run = shape.text_frame.paragraphs[0].runs[0]
    run.text = new_text


def fill_bulleted_list(shape, items: list):
    """Ajusta um textbox com bullets para ter exatamente len(items) parágrafos,
    clonando o XML do primeiro parágrafo (preserva fonte/cor/bullet)."""
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    if not items:
        items = ["-"]

    template_p = paragraphs[0]._p
    txBody = template_p.getparent()

    # remove todos os parágrafos existentes
    for p in list(paragraphs):
        txBody.remove(p._p)

    for text in items:
        new_p = copy.deepcopy(template_p)
        # ajusta o texto do run (mantém só o primeiro run, remove os demais)
        a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        runs = new_p.findall(f"{a_ns}r")
        for extra in runs[1:]:
            new_p.remove(extra)
        r = new_p.find(f"{a_ns}r")
        t = r.find(f"{a_ns}t")
        t.text = text
        txBody.append(new_p)


def fill_table(table_shape, entregas: list, max_bottom_emu: int = None):
    tbl = table_shape.table
    tbl_elem = tbl._tbl
    rows = list(tbl.rows)
    header_row = rows[0]
    data_rows = rows[1:]
    template_tr = data_rows[0]._tr

    n_needed = max(len(entregas), 1)
    n_have = len(data_rows)

    # adiciona linhas clonando a última linha de dados
    while n_have < n_needed:
        new_tr = copy.deepcopy(data_rows[-1]._tr)
        tbl_elem.append(new_tr)
        tbl = table_shape.table  # recarrega
        data_rows = list(tbl.rows)[1:]
        n_have = len(data_rows)

    # remove linhas em excesso
    tbl = table_shape.table
    data_rows = list(tbl.rows)[1:]
    while len(data_rows) > n_needed:
        tr = data_rows[-1]._tr
        tbl_elem.remove(tr)
        data_rows.pop()

    if not entregas:
        entregas = [{"nome": "-", "responsavel": "-", "status": "-", "obs": "-"}]

    tbl = table_shape.table
    data_rows = list(tbl.rows)[1:]
    for row, item in zip(data_rows, entregas):
        cells = row.cells
        cells[0].text_frame.paragraphs[0].runs[0].text = item["nome"] or "-"
        cells[1].text_frame.paragraphs[0].runs[0].text = item["responsavel"] or "-"
        cells[3].text_frame.paragraphs[0].runs[0].text = item["obs"] or "-"

        status_cell = cells[2]
        status_cell.text_frame.paragraphs[0].runs[0].text = item["status"] or "-"
        fg, bg = status_colors(item["status"])
        status_cell.fill.solid()
        status_cell.fill.fore_color.rgb = bg
        run = status_cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = fg
        run.font.bold = True

    # Ajusta a altura das linhas para caber no espaço disponível acima dos
    # cards de baixo, caso o número de entregas seja diferente de 5.
    if max_bottom_emu is not None:
        total_rows = 1 + len(data_rows)
        available = max_bottom_emu - table_shape.top
        default_row_h = Inches(0.475)
        min_row_h = Inches(0.32)
        row_h = min(default_row_h, max(min_row_h, available // total_rows))
        for row in tbl.rows:
            row.height = int(row_h)
        if row_h == min_row_h and available // total_rows < min_row_h:
            print(
                "Aviso: muitas entregas para o espaço do template — "
                "a tabela pode ultrapassar a área reservada. Considere "
                "resumir observações ou dividir em duas semanas."
            )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(description="Preenche o onepage semanal a partir de um .txt")
    parser.add_argument("txt_file", help="Arquivo de texto com os dados da semana")
    parser.add_argument(
        "--template",
        default="Onepage_Entregas_Semanais.pptx",
        help="Arquivo .pptx modelo (default: Onepage_Entregas_Semanais.pptx)",
    )
    parser.add_argument(
        "--output",
        default=None,
        help="Arquivo .pptx de saída (default: Onepage_Entregas_Semanais_preenchido.pptx)",
    )
    args = parser.parse_args()

    txt_path = Path(args.txt_file)
    template_path = Path(args.template)
    output_path = Path(args.output) if args.output else Path(
        "Onepage_Entregas_Semanais_preenchido.pptx"
    )

    if not txt_path.exists():
        sys.exit(f"Arquivo de dados não encontrado: {txt_path}")
    if not template_path.exists():
        sys.exit(f"Template .pptx não encontrado: {template_path}")

    data = parse_data(read_text_file(txt_path))

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    shapes = slide.shapes

    set_single_run_text(shapes[IDX_EQUIPE], f"Equipe: {data['equipe'] or '[Nome da Equipe]'}")
    set_single_run_text(
        shapes[IDX_PERIODO], f"Período: {data['periodo'] or '[DD/MM] a [DD/MM/AAAA]'}"
    )
    set_single_run_text(
        shapes[IDX_RESPONSAVEL],
        f"Responsável pelo envio: {data['responsavel'] or '[Nome]'}",
    )

    card_top_emu = shapes[12].top  # topo do card "Destaques da semana"
    fill_table(shapes[IDX_TABELA], data["entregas"], max_bottom_emu=card_top_emu - Inches(0.15))
    fill_bulleted_list(shapes[IDX_DESTAQUES], data["destaques"])
    fill_bulleted_list(shapes[IDX_RISCOS], data["riscos"])

    prs.save(str(output_path))
    print(f"Arquivo gerado: {output_path}")


if __name__ == "__main__":
    main()